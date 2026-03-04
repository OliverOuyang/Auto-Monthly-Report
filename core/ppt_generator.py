# -*- coding: utf-8 -*-
"""PPT 幻灯片组装"""

from pptx import Presentation
from pptx.util import Pt, Cm, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path


def create_presentation() -> Presentation:
    """创建标准宽屏 PPT。"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs


def add_slide(prs: Presentation, meta: dict, styles: dict,
              chart_path: str | Path, lines: list[str]) -> None:
    """
    添加一页幻灯片。
    """
    global_styles = styles.get('global', {})
    title_color_hex = global_styles.get('title_color', '#8B4513')
    font_family = global_styles.get('font_family', 'Microsoft YaHei')

    # 解析 hex color
    tc = _hex_to_rgb(title_color_hex)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # ── Title ──
    slide_title = meta.get('slide_title', '')
    tb = slide.shapes.add_textbox(Cm(2), Cm(0.8), Cm(28), Cm(1.4))
    p = tb.text_frame.paragraphs[0]
    p.text = slide_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = tc
    p.font.name = font_family

    # ── Golden underline ──
    ul = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(2), Cm(2.45), Cm(30), Cm(0.12))
    ul.fill.solid()
    ul.fill.fore_color.rgb = RGBColor(0xDA, 0xA5, 0x20)
    ul.line.fill.background()

    # ── Summary text ──
    tb2 = slide.shapes.add_textbox(Cm(2), Cm(3.0), Cm(29), Cm(3.0))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p2 = tf2.paragraphs[0]
        else:
            p2 = tf2.add_paragraph()
        p2.text = line
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p2.font.name = font_family
        p2.space_after = Pt(4)
        p2.line_spacing = Pt(24)

    # ── Chart image ──
    slide.shapes.add_picture(str(chart_path), Cm(0.8), Cm(5.8), Cm(32), Cm(12.5))

    # ── Red square decoration (top right) ──
    sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(31.5), Cm(1.0), Cm(1.2), Cm(1.2))
    sq.fill.solid()
    sq.fill.fore_color.rgb = RGBColor(0xE8, 0x4C, 0x3D)
    sq.line.fill.background()

    # ── Bottom gradient bar ──
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(18.2), Cm(33.86), Cm(0.8))
    bf = bar.fill
    bf.gradient()
    bf.gradient_stops[0].color.rgb = RGBColor(0xC0, 0x39, 0x2B)
    bf.gradient_stops[0].position = 0.0
    bf.gradient_stops[1].color.rgb = RGBColor(0xFA, 0xDB, 0xD8)
    bf.gradient_stops[1].position = 1.0
    bar.line.fill.background()


def save(prs: Presentation, output_path: str | Path) -> Path:
    """保存 PPT 文件。"""
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def _hex_to_rgb(hex_str: str) -> RGBColor:
    """'#8B4513' -> RGBColor(0x8B, 0x45, 0x13)"""
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
