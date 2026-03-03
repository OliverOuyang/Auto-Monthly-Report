# -*- coding: utf-8 -*-
"""
月报自动化 - 首借交易图表生成
从源数据生成透视表，输出 HTML (ECharts) 和 PPT
"""

import pandas as pd
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
from matplotlib.font_manager import FontProperties
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path
import json

# ── Paths ──
BASE = Path(r"C:\Users\Oliver\Desktop\数禾工作\16_AI项目\月报自动化")
DATA_FILE = BASE / "Data" / "2月月报数据_大盘交易aiV2.xlsx"
EXPORT_HTML = BASE / "export" / "HTML"
EXPORT_PPT = BASE / "export" / "PPT"
EXPORT_HTML.mkdir(parents=True, exist_ok=True)
EXPORT_PPT.mkdir(parents=True, exist_ok=True)

# ── Colors (matching reference chart) ──
COLORS = {
    '当月首登M0': '#4472C4',
    '存量首登M0': '#ED7D31',
    '初审M1+':   '#A5A5A5',
    '非初审':     '#FFC000',
    'API回流':    '#5B9BD5',
}
LINE_COLOR = '#FF0000'
CATEGORIES = ['当月首登M0', '存量首登M0', '初审M1+', '非初审', 'API回流']

# ── 1. Data Processing ──
print("[1/4] Reading raw data and building pivot table...")

raw = pd.read_excel(DATA_FILE, sheet_name='首借交易源数据')
raw['month'] = pd.to_datetime(raw['month'])
raw['loan_principal_amount'] = pd.to_numeric(raw['loan_principal_amount'], errors='coerce').fillna(0)

# Filter out "其他" user_group
raw_filtered = raw[raw['user_group'] != '其他'].copy()

# Merge 非初审 categories
raw_filtered['user_group2'] = raw_filtered['user_group'].replace({
    '非初审-重申': '非初审',
    '非初审-重审及其他': '非初审',
})

# Pivot: month × user_group2, sum loan_principal_amount
pivot = raw_filtered.pivot_table(
    index='month',
    columns='user_group2',
    values='loan_principal_amount',
    aggfunc='sum',
    fill_value=0
)
pivot = pivot.sort_index()

# Exclude current incomplete month (Mar-26 only has partial data)
REPORT_MONTH = pd.Timestamp('2026-02-01')  # 报告月份：2月
pivot = pivot.loc[pivot.index <= REPORT_MONTH]

pivot = pivot.reindex(columns=CATEGORIES, fill_value=0)
pivot['总计'] = pivot[CATEGORIES].sum(axis=1)

# Convert to 亿
pivot_yi = pivot / 1e8

print(f"  Months: {pivot.index.min().strftime('%Y-%m')} ~ {pivot.index.max().strftime('%Y-%m')}")
print(f"  Shape: {pivot.shape}")

# Month labels like "Jan-25"
month_labels = [d.strftime('%b-%y') for d in pivot.index]

# ── 2. Compute summary text ──
print("[2/4] Computing summary text...")

# Latest complete month = REPORT_MONTH
latest_month = REPORT_MONTH
prev_month = pivot.index[pivot.index < latest_month][-1] if len(pivot.index[pivot.index < latest_month]) > 0 else None

latest_total_yi = pivot_yi.loc[latest_month, '总计']
if prev_month is not None:
    prev_total_yi = pivot_yi.loc[prev_month, '总计']
    diff_yi = latest_total_yi - prev_total_yi
    pct_change = diff_yi / prev_total_yi * 100
    month_num = latest_month.month
    direction = "下降" if diff_yi < 0 else "上升"
    summary_text = (
        f"{month_num}月1-7评级首借交易（24h口径）总额{latest_total_yi:.2f}亿元，"
        f"环比{'减少' if diff_yi < 0 else '增加'}{abs(diff_yi):.1f}亿，"
        f"{direction}{abs(pct_change):.1f}%"
    )
else:
    month_num = latest_month.month
    summary_text = f"{month_num}月1-7评级首借交易（24h口径）总额{latest_total_yi:.2f}亿元"

print(f"  Summary: {summary_text}")

# ── 3. Generate HTML with ECharts ──
print("[3/4] Generating HTML...")


def format_yi(val):
    """Format value in 亿 with appropriate precision"""
    if abs(val) < 0.01:
        return ""
    if abs(val) >= 1:
        return f"{val:.2f}亿"
    else:
        return f"{val:.2f}亿"


# Prepare data for ECharts
echarts_data = {}
for cat in CATEGORIES:
    echarts_data[cat] = [round(float(v), 4) for v in pivot_yi[cat]]
echarts_data['总计'] = [round(float(v), 4) for v in pivot_yi['总计']]

# Data labels for each segment
label_data = {}
for cat in CATEGORIES:
    label_data[cat] = [format_yi(v) for v in pivot_yi[cat]]
label_data['总计'] = [format_yi(v) for v in pivot_yi['总计']]

html_content = f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>首借交易额-24H</title>
    <script src="https://cdn.jsdelivr.net/npm/echarts@5/dist/echarts.min.js"></script>
    <style>
        * {{ margin: 0; padding: 0; box-sizing: border-box; }}
        body {{
            font-family: 'Microsoft YaHei', 'SimHei', sans-serif;
            background: #fff;
            padding: 40px;
        }}
        .report-container {{
            max-width: 1200px;
            margin: 0 auto;
        }}
        .page-title {{
            font-size: 24px;
            font-weight: bold;
            color: #8B4513;
            margin-bottom: 4px;
        }}
        .title-underline {{
            height: 3px;
            background: linear-gradient(to right, #DAA520, #F0E68C);
            margin-bottom: 20px;
            width: 100%;
        }}
        .summary-text {{
            font-size: 16px;
            color: #333;
            line-height: 1.8;
            margin-bottom: 20px;
        }}
        #chart-container {{
            width: 100%;
            height: 520px;
        }}
        .footer-bar {{
            height: 8px;
            background: linear-gradient(to right, #C0392B, #F5B7B1, #FADBD8);
            margin-top: 30px;
            border-radius: 2px;
        }}
    </style>
</head>
<body>
    <div class="report-container">
        <div class="page-title">获客规模指标：1-7评级首借交易 By客群</div>
        <div class="title-underline"></div>
        <div class="summary-text">{summary_text}</div>
        <div id="chart-container"></div>
        <div class="footer-bar"></div>
    </div>
    <script>
        var chart = echarts.init(document.getElementById('chart-container'));
        var months = {json.dumps(month_labels, ensure_ascii=False)};
        var categories = {json.dumps(CATEGORIES, ensure_ascii=False)};
        var colors = {json.dumps(COLORS, ensure_ascii=False)};
        var lineColor = '{LINE_COLOR}';

        var data = {json.dumps(echarts_data, ensure_ascii=False)};
        var labels = {json.dumps(label_data, ensure_ascii=False)};

        var series = [];

        // Stacked bars
        categories.forEach(function(cat, idx) {{
            series.push({{
                name: cat,
                type: 'bar',
                stack: 'total',
                barWidth: '45%',
                itemStyle: {{ color: colors[cat] }},
                data: data[cat],
                label: {{
                    show: true,
                    position: 'inside',
                    formatter: function(params) {{
                        return labels[cat][params.dataIndex];
                    }},
                    fontSize: 11,
                    color: '#333',
                    fontWeight: 'bold'
                }}
            }});
        }});

        // Total line
        series.push({{
            name: '总计',
            type: 'line',
            data: data['总计'],
            symbol: 'circle',
            symbolSize: 6,
            lineStyle: {{
                color: lineColor,
                width: 3
            }},
            itemStyle: {{
                color: lineColor
            }},
            label: {{
                show: true,
                position: 'top',
                formatter: function(params) {{
                    return labels['总计'][params.dataIndex];
                }},
                fontSize: 12,
                color: lineColor,
                fontWeight: 'bold'
            }},
            z: 10
        }});

        var option = {{
            title: {{
                text: '首借交易额-24H',
                left: 'center',
                top: 0,
                textStyle: {{
                    fontSize: 16,
                    fontWeight: 'bold',
                    color: '#333'
                }}
            }},
            tooltip: {{
                trigger: 'axis',
                axisPointer: {{ type: 'shadow' }},
                formatter: function(params) {{
                    var html = '<b>' + params[0].axisValue + '</b><br/>';
                    var total = 0;
                    params.forEach(function(p) {{
                        if (p.seriesName !== '总计') {{
                            html += p.marker + ' ' + p.seriesName + ': ' + p.value.toFixed(2) + '亿<br/>';
                            total += p.value;
                        }}
                    }});
                    html += '<b>总计: ' + total.toFixed(2) + '亿</b>';
                    return html;
                }}
            }},
            legend: {{
                bottom: 0,
                data: categories.concat(['总计']),
                itemWidth: 14,
                itemHeight: 10,
                textStyle: {{ fontSize: 12 }}
            }},
            grid: {{
                left: '3%',
                right: '3%',
                top: 50,
                bottom: 50,
                containLabel: true
            }},
            xAxis: {{
                type: 'category',
                data: months,
                axisLabel: {{
                    fontSize: 11,
                    color: '#666'
                }},
                axisTick: {{ alignWithLabel: true }}
            }},
            yAxis: {{
                type: 'value',
                axisLabel: {{
                    formatter: function(v) {{ return v.toFixed(0) + '亿'; }},
                    fontSize: 11,
                    color: '#666'
                }},
                splitLine: {{
                    lineStyle: {{ color: '#eee' }}
                }}
            }},
            series: series
        }};

        chart.setOption(option);
        window.addEventListener('resize', function() {{ chart.resize(); }});
    </script>
</body>
</html>
"""

html_path = EXPORT_HTML / "首借交易额_24H.html"
with open(html_path, 'w', encoding='utf-8') as f:
    f.write(html_content)
print(f"  HTML saved: {html_path}")

# ── 4. Generate PPT ──
print("[4/4] Generating PPT...")

# First generate chart image with matplotlib
font_prop = FontProperties(fname=r'C:\Windows\Fonts\msyh.ttc')
plt.rcParams['font.sans-serif'] = ['Microsoft YaHei']
plt.rcParams['axes.unicode_minus'] = False

fig, ax = plt.subplots(figsize=(14, 6.5))

x = np.arange(len(month_labels))
bar_width = 0.55

# Draw stacked bars
bottoms = np.zeros(len(month_labels))
bars_dict = {}
for cat in CATEGORIES:
    values = pivot_yi[cat].values.astype(float)
    bars = ax.bar(x, values, bar_width, bottom=bottoms, color=COLORS[cat], label=cat, edgecolor='white', linewidth=0.3)
    bars_dict[cat] = (values, bottoms.copy())
    # Data labels inside bars
    for i, (v, b) in enumerate(zip(values, bottoms)):
        if v >= 0.3:  # Only label segments >= 0.3亿
            label_y = b + v / 2
            ax.text(i, label_y, f'{v:.2f}亿', ha='center', va='center',
                    fontsize=7, fontweight='bold', color='#333', fontproperties=font_prop)
    bottoms += values

# Total line
totals = pivot_yi['总计'].values.astype(float)
ax.plot(x, totals, 'o-', color=LINE_COLOR, linewidth=2.5, markersize=4, zorder=5, label='总计')
for i, v in enumerate(totals):
    ax.text(i, v + 0.3, f'{v:.2f}亿', ha='center', va='bottom',
            fontsize=8, fontweight='bold', color=LINE_COLOR, fontproperties=font_prop)

# Styling
ax.set_title('首借交易额-24H', fontsize=14, fontweight='bold', fontproperties=font_prop, pad=15)
ax.set_xticks(x)
ax.set_xticklabels(month_labels, fontsize=9)
ax.set_ylabel('')
ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v, _: ''))
ax.spines['top'].set_visible(False)
ax.spines['right'].set_visible(False)
ax.spines['left'].set_visible(False)
ax.tick_params(left=False)
ax.grid(axis='y', alpha=0.15)

# Legend at bottom
handles, labels_leg = ax.get_legend_handles_labels()
ax.legend(handles, labels_leg, loc='upper center', bbox_to_anchor=(0.5, -0.08),
          ncol=6, fontsize=9, frameon=False, prop=font_prop)

plt.tight_layout()
chart_img_path = EXPORT_PPT / "chart_temp.png"
fig.savefig(chart_img_path, dpi=200, bbox_inches='tight', facecolor='white')
plt.close()
print(f"  Chart image saved: {chart_img_path}")

# Build PPT
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

slide_layout = prs.slide_layouts[6]  # Blank
slide = prs.slides.add_slide(slide_layout)

# --- Background ---
from pptx.oxml.ns import qn
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# --- Page Title ---
from pptx.util import Cm
txBox = slide.shapes.add_textbox(Cm(1.5), Cm(0.8), Cm(30), Cm(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "获客规模指标：1-7评级首借交易 By客群"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = RGBColor(0x8B, 0x45, 0x13)
p.font.name = 'Microsoft YaHei'

# --- Title underline (golden) ---
from pptx.enum.shapes import MSO_SHAPE
line_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Cm(1.5), Cm(2.5), Cm(31), Cm(0.12)
)
line_shape.fill.solid()
line_shape.fill.fore_color.rgb = RGBColor(0xDA, 0xA5, 0x20)
line_shape.line.fill.background()

# --- Summary Text ---
txBox2 = slide.shapes.add_textbox(Cm(1.5), Cm(3.0), Cm(30), Cm(1.5))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = summary_text
p2.font.size = Pt(14)
p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
p2.font.name = 'Microsoft YaHei'
p2.line_spacing = Pt(22)

# --- Chart Image ---
slide.shapes.add_picture(
    str(chart_img_path), Cm(0.5), Cm(4.8), Cm(32), Cm(13.5)
)

# --- Bottom gradient bar ---
bar_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Cm(0), Cm(18.2), Cm(33.86), Cm(0.8)
)
bar_fill = bar_shape.fill
bar_fill.gradient()
bar_fill.gradient_stops[0].color.rgb = RGBColor(0xC0, 0x39, 0x2B)
bar_fill.gradient_stops[0].position = 0.0
bar_fill.gradient_stops[1].color.rgb = RGBColor(0xFA, 0xDB, 0xD8)
bar_fill.gradient_stops[1].position = 1.0
bar_shape.line.fill.background()

# --- Red square decoration (top right corner, matching PDF) ---
sq = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Cm(31.5), Cm(1.0), Cm(1.2), Cm(1.2)
)
sq.fill.solid()
sq.fill.fore_color.rgb = RGBColor(0xE8, 0x4C, 0x3D)
sq.line.fill.background()

ppt_path = EXPORT_PPT / "首借交易额_24H.pptx"
prs.save(str(ppt_path))
print(f"  PPT saved: {ppt_path}")

# Clean up temp chart image
# chart_img_path.unlink()  # Keep for debugging

print("\nDone. Output files:")
print(f"  HTML: {html_path}")
print(f"  PPT:  {ppt_path}")
